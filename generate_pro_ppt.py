#!/usr/bin/env python3
"""
Professional PPT matching Chitkara University official template
- White background
- Red top bar + Red bottom bar
- Chitkara logo top-right
- Times New Roman font throughout
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
from lxml import etree
from datetime import datetime

# ── Palette ────────────────────────────────────────────────────────────────────
RED       = RGBColor(0xCC, 0x00, 0x00)
BLACK     = RGBColor(0x00, 0x00, 0x00)
WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
DARK_GRAY = RGBColor(0x33, 0x33, 0x33)
MID_GRAY  = RGBColor(0x55, 0x55, 0x55)
LIGHT     = RGBColor(0xF5, 0xF5, 0xF5)
BORDER    = RGBColor(0xDD, 0xDD, 0xDD)
LIGHT_RED = RGBColor(0xFF, 0xEE, 0xEE)

FONT = "Times New Roman"

# Slide size (4:3 to match the example image more closely, but we'll do 16:9)
W = Inches(10)
H = Inches(7.5)

# ── Project info ───────────────────────────────────────────────────────────────
TITLE   = "An Intelligent DCNN-Based Framework for Anomaly Detection and Performance Degradation Analysis"
SUBTITLE = "Video Forgery Detection Using Machine Learning"
SCHOOL  = "Chitkara University, Rajpura, Punjab"
DEPT    = "Department of Computer Science and Engineering"
COURSE  = "IOHE (22CS422)"
GUIDE   = "Dr. Rajat Takkar"
TEAM    = [
    ("Gaganveer Singh", "2210994783"),
    ("Gunjan Mehta",    "2210991590"),
    ("Vidur Sharma",    "2210992524"),
]

# ══════════════════════════════════════════════════════════════════════════════
#  Helpers
# ══════════════════════════════════════════════════════════════════════════════

def add_rect(slide, left, top, width, height, fill=None, line=None, lw=0):
    sh = slide.shapes.add_shape(1, left, top, width, height)
    sh.line.fill.background()
    if fill:
        sh.fill.solid()
        sh.fill.fore_color.rgb = fill
    else:
        sh.fill.background()
    if line and lw:
        sh.line.color.rgb = line
        sh.line.width = Pt(lw)
    else:
        sh.line.fill.background()
    return sh


def add_text(slide, left, top, width, height,
             text, size=14, bold=False, italic=False,
             color=BLACK, align=PP_ALIGN.LEFT,
             anchor=MSO_ANCHOR.TOP, wrap=True):
    txb = slide.shapes.add_textbox(left, top, width, height)
    tf  = txb.text_frame
    tf.word_wrap = wrap
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name  = FONT
    run.font.size  = Pt(size)
    run.font.bold  = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return txb, tf


def add_para(tf, text, size=12, bold=False, italic=False,
             color=BLACK, align=PP_ALIGN.LEFT, space_before=4):
    p = tf.add_paragraph()
    p.alignment = align
    p.space_before = Pt(space_before)
    run = p.add_run()
    run.text = text
    run.font.name   = FONT
    run.font.size   = Pt(size)
    run.font.bold   = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return p


def rgb_to_hex(c):
    return f'{c[0]:02X}{c[1]:02X}{c[2]:02X}'

def cell_style(cell, fill_rgb=None, font_color=BLACK,
               size=11, bold=False, align=PP_ALIGN.LEFT):
    tf = cell.text_frame
    for p in tf.paragraphs:
        p.alignment = align
        for run in p.runs:
            run.font.name  = FONT
            run.font.size  = Pt(size)
            run.font.bold  = bold
            run.font.color.rgb = font_color
    if fill_rgb:
        tc   = cell._tc
        tcPr = tc.get_or_add_tcPr()
        for old in tcPr.findall(qn('a:solidFill')):
            tcPr.remove(old)
        solidFill = etree.SubElement(tcPr, qn('a:solidFill'))
        srgb = etree.SubElement(solidFill, qn('a:srgbClr'))
        r, g, b = fill_rgb[0], fill_rgb[1], fill_rgb[2]
        srgb.set('val', f'{r:02X}{g:02X}{b:02X}')


def make_table(slide, data, col_widths, row_height,
               left, top, hdr=True):
    rows = len(data)
    cols = len(data[0])
    tbl  = slide.shapes.add_table(
        rows, cols, left, top, sum(col_widths), row_height * rows
    ).table
    for c, w in enumerate(col_widths):
        tbl.columns[c].width = w
    for r in range(rows):
        tbl.rows[r].height = row_height
    for r, row in enumerate(data):
        for c, val in enumerate(row):
            cell = tbl.cell(r, c)
            cell.text = str(val)
            is_hdr = hdr and r == 0
            cell_style(
                cell,
                fill_rgb=RED if is_hdr else (LIGHT if r % 2 == 0 else WHITE),
                font_color=WHITE if is_hdr else DARK_GRAY,
                size=10 if is_hdr else 10,
                bold=is_hdr,
                align=PP_ALIGN.CENTER
            )
    return tbl


# ── Chrome: red bars top + bottom, logo top-right ─────────────────────────────
def chrome(slide, page_num=None):
    # Top red bar
    add_rect(slide, 0, 0, W, Inches(0.18), fill=RED)
    # Bottom red bar
    add_rect(slide, 0, H - Inches(0.18), W, Inches(0.18), fill=RED)

    # Logo box top-right  (white box with red "CU" badge style)
    logo_l = W - Inches(1.55)
    logo_t = Inches(0.18)
    logo_w = Inches(1.5)
    logo_h = Inches(0.72)
    add_rect(slide, logo_l, logo_t, logo_w, logo_h,
             fill=WHITE, line=BORDER, lw=0.5)
    # Red square inside logo
    add_rect(slide, logo_l + Inches(0.88), logo_t + Inches(0.08),
             Inches(0.52), Inches(0.56), fill=RED)
    # White "CU" mark on red square
    add_text(slide,
             logo_l + Inches(0.88), logo_t + Inches(0.08),
             Inches(0.52), Inches(0.56),
             "CU", size=14, bold=True, color=WHITE,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    # "CHITKARA" text
    add_text(slide, logo_l + Inches(0.05), logo_t + Inches(0.06),
             Inches(0.82), Inches(0.28),
             "CHITKARA", size=9, bold=True, color=BLACK,
             align=PP_ALIGN.LEFT)
    add_text(slide, logo_l + Inches(0.05), logo_t + Inches(0.34),
             Inches(0.82), Inches(0.26),
             "UNIVERSITY", size=8, bold=False, color=BLACK,
             align=PP_ALIGN.LEFT)

    # Page number bottom-right
    if page_num:
        add_text(slide, W - Inches(0.6), H - Inches(0.42),
                 Inches(0.5), Inches(0.28),
                 str(page_num), size=9, color=WHITE,
                 align=PP_ALIGN.CENTER)


def blank(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


# ══════════════════════════════════════════════════════════════════════════════
#  SLIDE 1 — Title
# ══════════════════════════════════════════════════════════════════════════════
def s1_title(prs):
    sl = blank(prs)
    add_rect(sl, 0, 0, W, H, fill=WHITE)
    chrome(sl)

    # Course header (bold, centered, top)
    add_text(sl, Inches(0.3), Inches(0.25), W - Inches(1.9), Inches(0.38),
             "Project Presentation of CO-OP Project at Industry (Module-2)  (22CS422)",
             size=13, bold=True, color=BLACK,
             align=PP_ALIGN.CENTER)

    add_text(sl, Inches(0.3), Inches(0.65), W - Inches(1.9), Inches(0.3),
             "On", size=13, color=BLACK, align=PP_ALIGN.CENTER)

    # Big project title
    add_text(sl, Inches(0.4), Inches(1.0), W - Inches(0.8), Inches(1.8),
             TITLE, size=28, bold=True, color=BLACK,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

    # Subtitle
    add_text(sl, Inches(0.4), Inches(2.85), W - Inches(0.8), Inches(0.45),
             SUBTITLE, size=14, italic=True, color=DARK_GRAY,
             align=PP_ALIGN.CENTER)

    # Team
    add_text(sl, Inches(0.4), Inches(3.45), W - Inches(0.8), Inches(0.3),
             TEAM[0][0] + "  |  " + TEAM[1][0] + "  |  " + TEAM[2][0],
             size=13, color=BLACK, align=PP_ALIGN.CENTER)
    add_text(sl, Inches(0.4), Inches(3.78), W - Inches(0.8), Inches(0.28),
             TEAM[0][1] + "  |  " + TEAM[1][1] + "  |  " + TEAM[2][1],
             size=12, color=MID_GRAY, align=PP_ALIGN.CENTER)

    add_text(sl, Inches(0.4), Inches(4.2), W - Inches(0.8), Inches(0.3),
             "Supervised By", size=13, color=BLACK, align=PP_ALIGN.CENTER)
    add_text(sl, Inches(0.4), Inches(4.5), W - Inches(0.8), Inches(0.3),
             GUIDE, size=13, color=BLACK, align=PP_ALIGN.CENTER)

    # Dept + Uni bottom
    add_text(sl, Inches(0.4), Inches(5.8), W - Inches(0.8), Inches(0.4),
             DEPT + ",", size=16, bold=True, color=BLACK,
             align=PP_ALIGN.CENTER)
    add_text(sl, Inches(0.4), Inches(6.2), W - Inches(0.8), Inches(0.4),
             SCHOOL, size=16, bold=True, color=BLACK,
             align=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════════════════════════════════
#  SLIDE 2 — Problem Statement
# ══════════════════════════════════════════════════════════════════════════════
def s2_problem(prs):
    sl = blank(prs)
    add_rect(sl, 0, 0, W, H, fill=WHITE)
    chrome(sl, 2)

    add_text(sl, Inches(0.3), Inches(0.22), W - Inches(1.9), Inches(0.42),
             "Problem Statement", size=20, bold=True, color=BLACK,
             align=PP_ALIGN.CENTER)

    # Horizontal red divider
    add_rect(sl, Inches(0.3), Inches(0.7), W - Inches(0.6), Inches(0.04), fill=RED)

    problems = [
        ("1.", "Digital Forgery Epidemic",
         "With the rise of video editing tools and deepfakes, malicious actors easily create forged videos that appear completely authentic, posing severe risks to truth and justice."),
        ("2.", "Failures of Existing Systems",
         "Current single-feature methods (keypoint matching, basic CNNs) achieve only 70-88% accuracy and cannot distinguish between active (frame insertion) and passive (compression) forgeries."),
        ("3.", "High-Stakes Consequences",
         "Forged videos are being used in criminal cases, fake news, surveillance fraud, and intellectual property theft — domains where a single detection failure has serious consequences."),
        ("4.", "No Comprehensive Solution",
         "No existing system combines multi-domain features (temporal + frequency + texture + statistical) with deep learning in a single automated, interpretable pipeline."),
    ]

    for i, (num, heading, desc) in enumerate(problems):
        y = Inches(0.85) + i * Inches(1.52)
        # Number badge
        add_rect(sl, Inches(0.3), y + Inches(0.1), Inches(0.42), Inches(0.42), fill=RED)
        add_text(sl, Inches(0.3), y + Inches(0.1), Inches(0.42), Inches(0.42),
                 num, size=13, bold=True, color=WHITE,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        add_text(sl, Inches(0.85), y + Inches(0.1), W - Inches(1.5), Inches(0.35),
                 heading, size=13, bold=True, color=BLACK)
        add_text(sl, Inches(0.85), y + Inches(0.48), W - Inches(1.5), Inches(0.95),
                 desc, size=11, color=MID_GRAY, wrap=True)


# ══════════════════════════════════════════════════════════════════════════════
#  SLIDE 3 — Existing Solutions (Gap Analysis)
# ══════════════════════════════════════════════════════════════════════════════
def s3_existing(prs):
    sl = blank(prs)
    add_rect(sl, 0, 0, W, H, fill=WHITE)
    chrome(sl, 3)

    add_text(sl, Inches(0.3), Inches(0.22), W - Inches(1.9), Inches(0.42),
             "Existing Solutions & Research Gap", size=20, bold=True, color=BLACK,
             align=PP_ALIGN.CENTER)
    add_rect(sl, Inches(0.3), Inches(0.7), W - Inches(0.6), Inches(0.04), fill=RED)

    # Comparison table
    add_text(sl, Inches(0.3), Inches(0.82), Inches(4), Inches(0.32),
             "Approach Comparison", size=13, bold=True, color=BLACK)

    tbl_data = [
        ["Method",                "Accuracy", "Forgery Types", "Limitation"],
        ["Keypoint (SIFT/SURF)",  "70-80%",   "Copy-Move",     "Fails on blur/compression"],
        ["Statistical Methods",   "75-85%",   "Passive only",  "Limited feature scope"],
        ["Single CNN",            "80-88%",   "Generic",       "Not interpretable"],
        ["ResNet50 Baseline",     "~79%",      "Generic",       "Raw pixels only"],
        ["OUR SYSTEM",            "85-95%",   "Active+Passive","Comprehensive (7 features)"],
    ]
    make_table(sl, tbl_data,
               col_widths=[Inches(2.2), Inches(1.0), Inches(1.5), Inches(2.5)],
               row_height=Inches(0.42),
               left=Inches(0.3), top=Inches(1.18))

    # Gap box
    add_rect(sl, Inches(0.3), Inches(4.28), W - Inches(0.6), Inches(0.04), fill=RED)
    add_rect(sl, Inches(0.3), Inches(4.34), W - Inches(0.6), Inches(1.2),
             fill=LIGHT_RED, line=RED, lw=1)
    add_text(sl, Inches(0.5), Inches(4.4), W - Inches(1.0), Inches(0.35),
             "Research Gap — Why Existing Solutions Are Not Enough:",
             size=13, bold=True, color=RED)
    add_text(sl, Inches(0.5), Inches(4.78), W - Inches(1.0), Inches(0.7),
             "No existing system simultaneously combines temporal, frequency-domain, texture, "
             "shape, structural, and statistical analysis in a single unified pipeline. "
             "Existing methods either lack coverage (single feature) or interpretability (pure DL). "
             "Our work fills this gap with a patent-eligible 7-feature fusion approach.",
             size=11, color=DARK_GRAY, wrap=True)

    # Three gap bullets at bottom
    gaps = [
        "Single feature coverage",
        "No active+passive distinction",
        "Poor cross-dataset generalization",
        "Black-box — not interpretable",
    ]
    for i, g in enumerate(gaps):
        x = Inches(0.3) + (i % 2) * Inches(4.7)
        y = Inches(5.65) + (i // 2) * Inches(0.45)
        add_rect(sl, x, y + Inches(0.1), Inches(0.18), Inches(0.18), fill=RED)
        add_text(sl, x + Inches(0.28), y, Inches(4.3), Inches(0.4),
                 g, size=11, bold=True, color=DARK_GRAY)


# ══════════════════════════════════════════════════════════════════════════════
#  SLIDE 4 — Our Solution
# ══════════════════════════════════════════════════════════════════════════════
def s4_solution(prs):
    sl = blank(prs)
    add_rect(sl, 0, 0, W, H, fill=WHITE)
    chrome(sl, 4)

    add_text(sl, Inches(0.3), Inches(0.22), W - Inches(1.9), Inches(0.42),
             "Our Solution — Multi-Feature Fusion Framework", size=20, bold=True,
             color=BLACK, align=PP_ALIGN.CENTER)
    add_rect(sl, Inches(0.3), Inches(0.7), W - Inches(0.6), Inches(0.04), fill=RED)

    # Tagline
    add_text(sl, Inches(0.3), Inches(0.8), W - Inches(0.6), Inches(0.38),
             "Hybrid Approach:  Traditional Computer Vision  +  Deep Learning (ResNet50 Transfer Learning)",
             size=12, bold=True, italic=True, color=MID_GRAY,
             align=PP_ALIGN.CENTER)

    # 7 feature boxes
    features = [
        ("Frame\nDifference",  "Temporal\nInconsistency"),
        ("DCT\nTransform",     "Compression\nArtifacts"),
        ("DWT\n(Haar)",        "Texture\nAnalysis"),
        ("LBP\nPatterns",      "Local Texture\nPatterns"),
        ("Binarization",       "Shape &\nBoundary"),
        ("Morphology",         "Structural\nAnomalies"),
        ("Eigen Vectors\n(PCA)","Statistical\nDeviations"),
    ]
    bw = Inches(1.3)
    bh = Inches(1.35)
    gap = Inches(0.05)
    sx  = Inches(0.28)
    sy  = Inches(1.32)

    for i, (name, desc) in enumerate(features):
        x = sx + i * (bw + gap)
        # Red header cap
        add_rect(sl, x, sy, bw, Inches(0.32), fill=RED)
        add_text(sl, x, sy, bw, Inches(0.32),
                 str(i+1), size=12, bold=True, color=WHITE,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        # White body with border
        add_rect(sl, x, sy + Inches(0.32), bw, bh - Inches(0.32),
                 fill=WHITE, line=BORDER, lw=0.8)
        add_text(sl, x + Inches(0.05), sy + Inches(0.36), bw - Inches(0.1), Inches(0.52),
                 name, size=10, bold=True, color=BLACK,
                 align=PP_ALIGN.CENTER, wrap=True)
        add_text(sl, x + Inches(0.05), sy + Inches(0.88), bw - Inches(0.1), Inches(0.42),
                 desc, size=9, color=MID_GRAY,
                 align=PP_ALIGN.CENTER, wrap=True)

    # Arrow / pipeline bar
    add_rect(sl, Inches(0.28), Inches(2.78), W - Inches(0.56), Inches(0.04), fill=RED)

    add_text(sl, Inches(0.28), Inches(2.86), W - Inches(0.56), Inches(0.35),
             "7 Features  -->  12-Channel Tensor (240x320x12)  -->  ResNet50 Backbone (Frozen)  -->  Dense Classifier  -->  FORGED / AUTHENTIC",
             size=11, bold=True, color=BLACK, align=PP_ALIGN.CENTER)

    add_rect(sl, Inches(0.28), Inches(3.24), W - Inches(0.56), Inches(0.04), fill=RED)

    # KPI boxes
    kpis = [
        ("85 - 95%",  "Accuracy on Real Dataset"),
        ("+16%",      "Improvement vs Baseline"),
        ("24.8M",     "Total Model Parameters"),
        ("5-10 sec",  "Inference per Video"),
        ("7",         "Feature Extractors"),
        ("2,500+",    "Lines of Code"),
    ]
    kw = Inches(1.58)
    kh = Inches(0.88)
    ky = Inches(3.38)
    for i, (val, lbl) in enumerate(kpis):
        kx = Inches(0.28) + i * (kw + Inches(0.04))
        add_rect(sl, kx, ky, kw, kh, fill=RED)
        add_text(sl, kx, ky + Inches(0.04), kw, Inches(0.44),
                 val, size=18, bold=True, color=WHITE,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        add_text(sl, kx, ky + Inches(0.48), kw, Inches(0.36),
                 lbl, size=8.5, color=WHITE,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, wrap=True)

    # Features note
    note = ("The 7 feature extractors operate in parallel and are fused into a single 12-channel tensor. "
            "This ensemble approach makes it extremely difficult for any forgery to evade all detectors simultaneously, "
            "resulting in 85-95% accuracy on real datasets — a 6-16% gain over single-feature baselines.")
    add_text(sl, Inches(0.3), Inches(4.42), W - Inches(0.6), Inches(0.85),
             note, size=10.5, italic=True, color=MID_GRAY, wrap=True,
             align=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════════════════════════════════
#  SLIDE 5 — System Architecture
# ══════════════════════════════════════════════════════════════════════════════
def s5_architecture(prs):
    sl = blank(prs)
    add_rect(sl, 0, 0, W, H, fill=WHITE)
    chrome(sl, 5)

    add_text(sl, Inches(0.3), Inches(0.22), W - Inches(1.9), Inches(0.42),
             "System Architecture & Processing Pipeline", size=20, bold=True,
             color=BLACK, align=PP_ALIGN.CENTER)
    add_rect(sl, Inches(0.3), Inches(0.7), W - Inches(0.6), Inches(0.04), fill=RED)

    # Pipeline stages
    stages = [
        ("INPUT\nVIDEO",      "MP4 / AVI\nMOV / MKV"),
        ("FRAME\nEXTRACT",    "OpenCV\n240x320px"),
        ("PRE-\nPROCESS",     "Grayscale\n+ Canny"),
        ("7 FEATURE\nEXTRACT","Parallel\nMethods"),
        ("FEATURE\nFUSION",   "12-Channel\nTensor"),
        ("RESNET50\nFROZEN",  "23.6M\nparams"),
        ("DENSE\nHEAD",       "512-256\n-128-1"),
        ("OUTPUT",            "FORGED /\nAUTHENTIC"),
    ]
    bw  = Inches(1.12)
    bh  = Inches(1.4)
    sx  = Inches(0.28)
    sy  = Inches(0.84)
    gap = Inches(0.07)

    for i, (title, detail) in enumerate(stages):
        x = sx + i * (bw + gap)
        add_rect(sl, x, sy, bw, Inches(0.42), fill=RED)
        add_text(sl, x, sy, bw, Inches(0.42),
                 title, size=9, bold=True, color=WHITE,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, wrap=True)
        add_rect(sl, x, sy + Inches(0.42), bw, bh - Inches(0.42),
                 fill=LIGHT, line=BORDER, lw=0.6)
        add_text(sl, x + Inches(0.04), sy + Inches(0.46), bw - Inches(0.08),
                 bh - Inches(0.54),
                 detail, size=9, color=MID_GRAY,
                 align=PP_ALIGN.CENTER, wrap=True)
        if i < len(stages) - 1:
            ax = x + bw + Inches(0.01)
            add_text(sl, ax, sy + Inches(0.55), Inches(0.06), Inches(0.3),
                     ">", size=11, bold=True, color=RED, align=PP_ALIGN.CENTER)

    add_rect(sl, Inches(0.28), Inches(2.34), W - Inches(0.56), Inches(0.04), fill=RED)

    # Two spec tables side by side
    add_text(sl, Inches(0.3), Inches(2.44), Inches(4.6), Inches(0.32),
             "Model Specifications", size=13, bold=True, color=BLACK)

    spec = [
        ["Parameter",            "Value"],
        ["Total Parameters",     "24,827,076"],
        ["Trainable Params",     "1,239,172  (5%)"],
        ["Frozen (ResNet50)",    "23,587,904  (95%)"],
        ["Model File Size",      "94.71 MB"],
        ["Input Tensor",         "240 x 320 x 12"],
        ["Output",               "0=Authentic / 1=Forged"],
    ]
    make_table(sl, spec,
               col_widths=[Inches(2.1), Inches(2.5)],
               row_height=Inches(0.4),
               left=Inches(0.3), top=Inches(2.8))

    add_text(sl, Inches(5.05), Inches(2.44), Inches(4.6), Inches(0.32),
             "Training Configuration", size=13, bold=True, color=BLACK)

    train = [
        ["Config",           "Value"],
        ["Optimizer",        "Adam  (lr=0.0001)"],
        ["Loss",             "Binary Crossentropy"],
        ["Batch Size",       "16"],
        ["Epochs",           "10 - 20"],
        ["Val Split",        "80% Train / 20% Val"],
        ["Training Time",    "10-15 minutes (CPU)"],
    ]
    make_table(sl, train,
               col_widths=[Inches(1.6), Inches(3.0)],
               row_height=Inches(0.4),
               left=Inches(5.05), top=Inches(2.8))


# ══════════════════════════════════════════════════════════════════════════════
#  SLIDE 6 — Demo & Screenshots
# ══════════════════════════════════════════════════════════════════════════════
def s6_demo(prs):
    sl = blank(prs)
    add_rect(sl, 0, 0, W, H, fill=WHITE)
    chrome(sl, 6)

    add_text(sl, Inches(0.3), Inches(0.22), W - Inches(1.9), Inches(0.42),
             "Demo & Workflow — How to Use the System", size=20, bold=True,
             color=BLACK, align=PP_ALIGN.CENTER)
    add_rect(sl, Inches(0.3), Inches(0.7), W - Inches(0.6), Inches(0.04), fill=RED)

    # Left: steps
    steps = [
        ("Step 1", "Install Dependencies",
         "pip install tensorflow opencv-python scikit-image PyWavelets scikit-learn"),
        ("Step 2", "Train the Model  (~10-15 min)",
         "python run_complete_pipeline.py\n>> Select option 1  >>  Type 'y' to train"),
        ("Step 3", "Predict a Video",
         "python predict_video.py  path/to/video.mp4"),
        ("Step 4", "Read the Output",
         "Video: test.mp4\nResult:  FORGED\nConfidence:  87.34%"),
    ]
    for i, (step, title, cmd) in enumerate(steps):
        y = Inches(0.84) + i * Inches(1.52)
        add_rect(sl, Inches(0.3), y, Inches(1.0), Inches(0.42), fill=RED)
        add_text(sl, Inches(0.3), y, Inches(1.0), Inches(0.42),
                 step, size=10, bold=True, color=WHITE,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        add_text(sl, Inches(1.42), y, Inches(4.2), Inches(0.38),
                 title, size=12, bold=True, color=BLACK)
        # Code-style box
        add_rect(sl, Inches(1.42), y + Inches(0.44), Inches(4.2), Inches(0.92),
                 fill=RGBColor(0xF0, 0xF0, 0xF0), line=BORDER, lw=0.5)
        add_text(sl, Inches(1.55), y + Inches(0.48), Inches(3.95), Inches(0.84),
                 cmd, size=9.5, color=DARK_GRAY, wrap=True)

    add_rect(sl, Inches(0.28), Inches(0.84), Inches(0.04), Inches(5.82), fill=RED)

    # Right: 12-channel feature diagram
    add_text(sl, Inches(5.9), Inches(0.84), Inches(3.8), Inches(0.38),
             "12-Channel Feature Tensor", size=13, bold=True, color=BLACK)
    add_rect(sl, Inches(5.9), Inches(1.26), Inches(3.8), Inches(0.04), fill=RED)

    channels = [
        ("Ch 0",    "Grayscale"),
        ("Ch 1",    "Canny Edges"),
        ("Ch 2",    "DCT"),
        ("Ch 3-6",  "DWT (4 sub-bands)"),
        ("Ch 7",    "LBP Texture"),
        ("Ch 8",    "Binarization"),
        ("Ch 9-10", "Morphology"),
        ("Ch 11",   "Eigen Vectors"),
    ]
    ch_data = [["Channel", "Feature"]] + list(channels)
    make_table(sl, ch_data,
               col_widths=[Inches(1.0), Inches(2.72)],
               row_height=Inches(0.4),
               left=Inches(5.9), top=Inches(1.34))

    # Output simulation box
    add_text(sl, Inches(5.9), Inches(5.0), Inches(3.8), Inches(0.35),
             "Sample Terminal Output:", size=12, bold=True, color=BLACK)
    add_rect(sl, Inches(5.9), Inches(5.38), Inches(3.8), Inches(1.28),
             fill=RGBColor(0x1E, 0x1E, 0x1E))

    txb = sl.shapes.add_textbox(Inches(6.0), Inches(5.44), Inches(3.6), Inches(1.14))
    tf  = txb.text_frame
    tf.word_wrap = True
    lines = [
        ("VIDEO FORGERY DETECTION",   WHITE,                 True,  9),
        ("================================", MID_GRAY,       False, 8),
        ("Video:       test_video.mp4",  RGBColor(0xAA,0xAA,0xAA), False, 9),
        ("Result:      FORGED",        RGBColor(0xFF,0x44,0x44), True,  10),
        ("Confidence:  87.34%",        RGBColor(0x00,0xEE,0x76), True,  10),
        ("================================", MID_GRAY,       False, 8),
    ]
    for j, (line, col, bld, sz) in enumerate(lines):
        p = tf.paragraphs[0] if j == 0 else tf.add_paragraph()
        run = p.add_run()
        run.text = line
        run.font.name  = "Courier New"
        run.font.size  = Pt(sz)
        run.font.bold  = bld
        run.font.color.rgb = col


# ══════════════════════════════════════════════════════════════════════════════
#  SLIDE 7 — Patent / IPR
# ══════════════════════════════════════════════════════════════════════════════
def s7_patent(prs):
    sl = blank(prs)
    add_rect(sl, 0, 0, W, H, fill=WHITE)
    chrome(sl, 7)

    add_text(sl, Inches(0.3), Inches(0.22), W - Inches(1.9), Inches(0.42),
             "Patent & IPR Opportunity", size=20, bold=True, color=BLACK,
             align=PP_ALIGN.CENTER)
    add_rect(sl, Inches(0.3), Inches(0.7), W - Inches(0.6), Inches(0.04), fill=RED)

    # Patent claims (left)
    add_text(sl, Inches(0.3), Inches(0.82), Inches(5.5), Inches(0.35),
             "Novel Patent Claims", size=13, bold=True, color=BLACK)

    claims = [
        "Multi-Feature Fusion: First system combining 7 complementary feature extractors (DCT + DWT + LBP + Morphology + PCA + Temporal + Binary) into a unified 12-channel tensor for video forensics.",
        "Hybrid CV + DL Architecture: Novel pairing of interpretable traditional computer vision with frozen ResNet50 transfer learning — providing both accuracy (85-95%) and explainability.",
        "Active + Passive Unified Detection: Single automated pipeline detecting both temporal-based active forgeries and compression-based passive forgeries simultaneously.",
        "Automated Forensic Pipeline: End-to-end system from raw video input to verdict with confidence score — no manual feature engineering or configuration required.",
    ]
    for i, claim in enumerate(claims):
        y = Inches(1.22) + i * Inches(1.38)
        add_rect(sl, Inches(0.3), y, Inches(0.38), Inches(0.38), fill=RED)
        add_text(sl, Inches(0.3), y, Inches(0.38), Inches(0.38),
                 str(i+1), size=13, bold=True, color=WHITE,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        add_text(sl, Inches(0.8), y, Inches(5.0), Inches(1.25),
                 claim, size=10, color=DARK_GRAY, wrap=True)

    add_rect(sl, Inches(5.62), Inches(0.82), Inches(0.04), Inches(5.82), fill=RED)

    # Right: Inventor details table
    add_text(sl, Inches(5.78), Inches(0.82), Inches(3.9), Inches(0.35),
             "Inventor Details (Chalkpad Filing)", size=13, bold=True, color=BLACK)

    ipr = [
        ["Name",             "Emp Code",   "Mobile"],
        ["Gaganveer Singh",  "2210994783", "8295055789"],
        ["Gunjan Mehta",     "2210991590", "9315636376"],
        ["Vidur Sharma",     "2210992524", "8570840245"],
    ]
    make_table(sl, ipr,
               col_widths=[Inches(2.0), Inches(1.2), Inches(1.4)],
               row_height=Inches(0.42),
               left=Inches(5.78), top=Inches(1.22))

    # Emails
    emails = [
        ["Email (University)"],
        ["gaganveer4783.be22@chitkara.edu.in"],
        ["gunjan1590.be22@chitkara.edu.in"],
        ["vidur2524.be22@chitkara.edu.in"],
    ]
    make_table(sl, emails,
               col_widths=[Inches(3.9)],
               row_height=Inches(0.42),
               left=Inches(5.78), top=Inches(2.92))

    # Market box
    add_rect(sl, Inches(5.78), Inches(4.6), Inches(3.9), Inches(0.04), fill=RED)
    add_text(sl, Inches(5.78), Inches(4.7), Inches(3.9), Inches(0.32),
             "Commercial Applications", size=12, bold=True, color=BLACK)
    apps = [
        "Digital Forensics Agencies",
        "Law Enforcement & Courts",
        "Media Platforms (YouTube, Meta)",
        "Government Security Systems",
        "Enterprise Content Moderation",
    ]
    for i, app in enumerate(apps):
        y = Inches(5.08) + i * Inches(0.42)
        add_rect(sl, Inches(5.78), y + Inches(0.1), Inches(0.2), Inches(0.2), fill=RED)
        add_text(sl, Inches(6.08), y, Inches(3.55), Inches(0.38),
                 app, size=11, color=DARK_GRAY)


# ══════════════════════════════════════════════════════════════════════════════
#  SLIDE 8 — Results & Impact
# ══════════════════════════════════════════════════════════════════════════════
def s8_results(prs):
    sl = blank(prs)
    add_rect(sl, 0, 0, W, H, fill=WHITE)
    chrome(sl, 8)

    add_text(sl, Inches(0.3), Inches(0.22), W - Inches(1.9), Inches(0.42),
             "Results, Impact & Conclusion", size=20, bold=True, color=BLACK,
             align=PP_ALIGN.CENTER)
    add_rect(sl, Inches(0.3), Inches(0.7), W - Inches(0.6), Inches(0.04), fill=RED)

    # KPI strip
    kpis = [
        ("85-95%",  "Test Accuracy\n(Real Dataset)"),
        ("+16%",    "Improvement\nvs Baseline"),
        ("90%+",    "Active Forgery\nDetection"),
        ("85%+",    "Passive Forgery\nDetection"),
        ("<10%",    "False Positive\nRate"),
    ]
    kw = Inches(1.85)
    for i, (val, lbl) in enumerate(kpis):
        kx = Inches(0.28) + i * (kw + Inches(0.05))
        add_rect(sl, kx, Inches(0.82), kw, Inches(1.02), fill=RED)
        add_text(sl, kx, Inches(0.84), kw, Inches(0.52),
                 val, size=22, bold=True, color=WHITE,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        add_text(sl, kx, Inches(1.34), kw, Inches(0.42),
                 lbl, size=8.5, color=WHITE,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, wrap=True)

    # Comparison table left
    add_text(sl, Inches(0.3), Inches(2.0), Inches(5.3), Inches(0.32),
             "Benchmark Comparison", size=13, bold=True, color=BLACK)
    cmp = [
        ["Method",              "Accuracy", "Interpretable"],
        ["Keypoint Matching",   "70-80%",   "Yes"],
        ["Statistical Methods", "75-85%",   "Partial"],
        ["Single CNN",          "80-88%",   "No"],
        ["ResNet50 Baseline",   "~79%",     "No"],
        ["OUR SYSTEM",          "85-95%",   "Yes"],
    ]
    make_table(sl, cmp,
               col_widths=[Inches(2.5), Inches(1.2), Inches(1.5)],
               row_height=Inches(0.4),
               left=Inches(0.3), top=Inches(2.36))

    # Feature accuracy right
    add_text(sl, Inches(5.55), Inches(2.0), Inches(4.15), Inches(0.32),
             "Per-Feature Detection Rates", size=13, bold=True, color=BLACK)
    feat = [
        ["Feature",         "Rate",   "Forgery Type"],
        ["Frame Difference","90%+",   "Active"],
        ["DCT",             "88%+",   "Passive"],
        ["DWT",             "87%+",   "Both"],
        ["LBP",             "85%+",   "Both"],
        ["Morphology",      "84%+",   "Both"],
        ["Eigen Vectors",   "85%+",   "Both"],
    ]
    make_table(sl, feat,
               col_widths=[Inches(1.85), Inches(0.9), Inches(1.3)],
               row_height=Inches(0.4),
               left=Inches(5.55), top=Inches(2.36))

    # Achievements list
    add_rect(sl, Inches(0.28), Inches(5.28), W - Inches(0.56), Inches(0.04), fill=RED)
    add_text(sl, Inches(0.28), Inches(5.36), W - Inches(0.56), Inches(0.35),
             "Key Achievements & Conclusion", size=13, bold=True, color=BLACK)

    achievements = [
        "7 complementary feature extraction methods covering all forgery domains",
        "85-95% accuracy on real REWIND dataset — 6-16% improvement over all baselines",
        "Fully automated one-command pipeline — train in 15 min, predict in 10 sec",
        "Patent-eligible novel multi-feature fusion architecture",
        "Production-ready 2,500+ line codebase — ready for deployment",
    ]
    for i, ach in enumerate(achievements):
        y = Inches(5.78) + i * Inches(0.28)
        add_rect(sl, Inches(0.3), y + Inches(0.06), Inches(0.18), Inches(0.18), fill=RED)
        add_text(sl, Inches(0.6), y, W - Inches(0.95), Inches(0.28),
                 ach, size=10.5, bold=(i == 0), color=DARK_GRAY)


# ══════════════════════════════════════════════════════════════════════════════
#  Main
# ══════════════════════════════════════════════════════════════════════════════
def build(out):
    prs = Presentation()
    prs.slide_width  = W
    prs.slide_height = H

    s1_title(prs)
    s2_problem(prs)
    s3_existing(prs)
    s4_solution(prs)
    s5_architecture(prs)
    s6_demo(prs)
    s7_patent(prs)
    s8_results(prs)

    prs.save(out)
    print("[OK] Saved:", out)


if __name__ == "__main__":
    out = r"C:\Users\GunjanMehta\Desktop\patent\Video_Forgery_Presentation.pptx"
    print("Generating professional Chitkara-style PPT...")
    build(out)
    print("Done!")
